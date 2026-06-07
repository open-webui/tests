"""Dependency contract: python-pptx (import name ``pptx``).

Open WebUI's ``retrieval/loaders/main.py`` has a ``PptxLoader`` — the
fallback PowerPoint loader used when ``unstructured`` is not installed. It
uses a narrow but exact slice of python-pptx to extract slide text for RAG:

    from pptx import Presentation
    prs = Presentation(self.file_path)
    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_texts.append(shape.text_frame.text)

So the contract that must hold:
  - ``Presentation(source)`` opens a .pptx from a path OR a file-like object
    and yields a presentation whose ``.slides`` is iterable;
  - each slide's ``.shapes`` is iterable;
  - every shape exposes a boolean ``.has_text_frame``;
  - a text-bearing shape exposes ``.text_frame.text`` returning the shape's
    text with paragraphs joined by newlines.

This module builds a presentation entirely IN MEMORY (via ``io.BytesIO`` —
no disk, no network), round-trips it through ``Presentation`` exactly as the
loader does, and asserts the extraction the loader performs. A python-pptx
bump that renamed ``has_text_frame`` / ``text_frame`` / ``slides`` /
``shapes`` would break PowerPoint ingestion silently; this fails loudly.

Pattern mirrors test_requests.py. Uses the ``depcheck`` fixture.
"""

from __future__ import annotations

import io

import pytest

pytestmark = pytest.mark.depcheck

IMPORT_NAME = "pptx"
DIST_NAME = "python-pptx"


# ---------------------------------------------------------------------------
# In-memory builder — produces a .pptx byte stream with known slide text,
# without touching the filesystem.
# ---------------------------------------------------------------------------


def _build_pptx_bytes(mod, slides_text):
    """Return a BytesIO of a .pptx whose slides carry the given text.

    `slides_text` is a list of lists: one list of paragraph strings per slide.
    Each slide gets a single textbox holding those paragraphs.
    """
    util = _util(mod)
    prs = mod.Presentation()
    blank = prs.slide_layouts[6]  # the blank layout (no placeholders)
    for paragraphs in slides_text:
        slide = prs.slides.add_slide(blank)
        one, six, two = util.Inches(1), util.Inches(6), util.Inches(2)
        tb = slide.shapes.add_textbox(one, one, six, two)
        tf = tb.text_frame
        tf.text = paragraphs[0]
        for extra in paragraphs[1:]:
            tf.add_paragraph().text = extra
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def _util(mod):
    import importlib

    return importlib.import_module("pptx.util")


def _extract_like_loader(mod, source):
    """Reproduce PptxLoader.load()'s extraction against an opened source."""
    prs = mod.Presentation(source)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_texts.append(shape.text_frame.text)
        if slide_texts:
            parts.append(f"Slide {i}:\n" + "\n".join(slide_texts))
    return "\n\n".join(parts)


# ---------------------------------------------------------------------------
# Import + version + surface
# ---------------------------------------------------------------------------


def test_import(depcheck):
    mod = depcheck.load(IMPORT_NAME)
    assert mod.__name__ == "pptx"


def test_version_reported(depcheck):
    assert depcheck.dist_version(DIST_NAME) is not None


def test_presentation_is_callable(depcheck):
    """`from pptx import Presentation` then `Presentation(source)` — the entry
    point must be a callable factory."""
    mod = depcheck.load(IMPORT_NAME)
    assert callable(mod.Presentation)


def test_presentation_accepts_no_arg(depcheck):
    """Presentation() with no arg yields a blank deck (used to build fixtures;
    the signature must keep the source argument optional)."""
    mod = depcheck.load(IMPORT_NAME)
    prs = mod.Presentation()
    assert prs is not None
    assert len(prs.slides) == 0


# ---------------------------------------------------------------------------
# slides / shapes / text_frame — the exact object graph the loader walks.
# ---------------------------------------------------------------------------


def test_slides_is_iterable_and_enumerable(depcheck):
    """The loader does `for i, slide in enumerate(prs.slides, 1)`. slides must
    be iterable and have a length."""
    mod = depcheck.load(IMPORT_NAME)
    buf = _build_pptx_bytes(mod, [["A"], ["B"], ["C"]])
    prs = mod.Presentation(buf)
    assert len(prs.slides) == 3
    collected = [s for s in prs.slides]
    assert len(collected) == 3


def test_shape_has_text_frame_flag(depcheck):
    """Every shape must expose a boolean has_text_frame; a textbox shape must
    report True."""
    mod = depcheck.load(IMPORT_NAME)
    buf = _build_pptx_bytes(mod, [["only slide text"]])
    prs = mod.Presentation(buf)
    slide = next(iter(prs.slides))
    shapes = list(slide.shapes)
    assert shapes, "expected at least one shape on the slide"
    for shape in shapes:
        assert isinstance(shape.has_text_frame, bool)
    assert any(s.has_text_frame for s in shapes), "no text-bearing shape found"


def test_text_frame_text_joins_paragraphs_with_newline(depcheck):
    """text_frame.text returns the shape's paragraphs joined by '\\n' — the
    loader relies on this to get multi-line slide text."""
    mod = depcheck.load(IMPORT_NAME)
    buf = _build_pptx_bytes(mod, [["First line", "Second line", "Third line"]])
    prs = mod.Presentation(buf)
    slide = next(iter(prs.slides))
    text_shapes = [s for s in slide.shapes if s.has_text_frame]
    assert text_shapes
    assert text_shapes[0].text_frame.text == "First line\nSecond line\nThird line"


# ---------------------------------------------------------------------------
# End-to-end: reproduce PptxLoader.load() output offline.
# ---------------------------------------------------------------------------


def test_loader_extraction_from_bytesio(depcheck):
    """Open a deck from a file-like object (BytesIO) exactly as the loader
    opens a path, and reproduce its concatenated 'Slide N:\\n...' output."""
    mod = depcheck.load(IMPORT_NAME)
    buf = _build_pptx_bytes(mod, [["Intro", "agenda"], ["Details"]])
    result = _extract_like_loader(mod, buf)
    assert result == "Slide 1:\nIntro\nagenda\n\nSlide 2:\nDetails"


def test_loader_extraction_from_real_path(depcheck, tmp_path):
    """The loader passes a filesystem path; confirm Presentation(path) works
    too (tmp_path is pytest's per-test temp dir — local, not network)."""
    mod = depcheck.load(IMPORT_NAME)
    buf = _build_pptx_bytes(mod, [["Hello"]])
    path = tmp_path / "deck.pptx"
    path.write_bytes(buf.getvalue())
    result = _extract_like_loader(mod, str(path))
    assert result == "Slide 1:\nHello"


def test_round_trip_preserves_text(depcheck):
    """Saving to BytesIO and reloading must preserve slide text (the loader
    only ever reads, but this pins the save/load fidelity the builder uses)."""
    mod = depcheck.load(IMPORT_NAME)
    buf = _build_pptx_bytes(mod, [["persisted text"]])
    prs = mod.Presentation(buf)
    texts = [s.text_frame.text for slide in prs.slides for s in slide.shapes if s.has_text_frame]
    assert "persisted text" in texts


def test_empty_presentation_yields_no_text(depcheck):
    """A deck with no slides produces empty extraction (the loader appends
    nothing) — guards the empty-input path."""
    mod = depcheck.load(IMPORT_NAME)
    prs = mod.Presentation()
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    assert _extract_like_loader(mod, buf) == ""


# ---------------------------------------------------------------------------
# util constructors the builder (and any layout-aware consumer) relies on.
# ---------------------------------------------------------------------------


def test_util_inches_available(depcheck):
    """pptx.util.Inches is used to position the textbox; it must remain a
    callable producing an EMU length."""
    depcheck.load(IMPORT_NAME)
    util = _util(depcheck.load(IMPORT_NAME))
    assert callable(util.Inches)
    assert int(util.Inches(1)) > 0
